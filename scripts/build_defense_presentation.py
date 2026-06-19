"""Build defense presentation from university template."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = Path(r"c:\Users\tbnamozova\Downloads\Научный проект.pptx")
OUTPUT = PROJECT_ROOT / "docs" / "Защита_MHSNet_Fraud_Detection.pptx"

TOPIC = (
  "Разработка метода выявления мошеннических транзакций "
  "с применением сети Хопфилда (MHSNet)"
)

SLIDE_CONTENT: dict[int, str] = {
  1: (
    f"Тема: {TOPIC}\n\n"
    "Студент: Намозова Томирис Борисовна, гр. [указать группу]\n\n"
    "Научный руководитель: [ФИО, учёная степень, должность]\n\n"
    "Научный консультант: [при наличии]"
  ),
  2: (
    "Мошеннические банковские транзакции приводят к прямым финансовым потерям "
    "и снижению доверия клиентов. В реальных данных доля fraud крайне мала "
    "(около 0,1%), а метки мошенничества часто поступают с задержкой (chargeback).\n\n"
    "Классические supervised-модели требуют большого объёма размеченных примеров fraud "
    "и плохо переносятся на новые схемы мошенничества (concept drift).\n\n"
    "Актуально исследовать unsupervised и гибридные методы, которые учатся на "
    "нормальных операциях и выявляют аномальные транзакции.\n\n"
    "Научность работы: адаптация пайплайна MHSNet (Kernel PCA → LOF → Modern Hopfield) "
    "к реальному банковскому датасету CaixaBank (~13 млн транзакций) с воспроизводимой "
    "экспериментальной оценкой."
  ),
  3: (
    "Цель исследования:\n"
    "Разработать и экспериментально исследовать метод выявления мошеннических "
    "транзакций на основе пайплайна MHSNet (Kernel PCA → LOF → сеть Хопфилда) "
    "на реальном датасете CaixaBank.\n\n"
    "Задачи:\n"
    "1. Проанализировать методы обнаружения мошеннических транзакций и аномалий "
    "в финансовых данных.\n"
    "2. Реализовать загрузку и предобработку датасета "
    "computingvictor/transactions-fraud-datasets (Kaggle).\n"
    "3. Реализовать пайплайн MHSNet: Kernel PCA, LOF, классическая сеть Хопфилда "
    "с Modern Hopfield Energy.\n"
    "4. Провести эксперименты и сравнение с baseline-моделями "
    "(Logistic Regression, Random Forest, Isolation Forest).\n"
    "5. Оценить качество по метрикам F1, ROC-AUC и Recall@FPR=5%."
  ),
  4: (
    "Zhao Y. MHSNet-SNN: Improved Outlier Detection for Fraud Detection in "
    "Financial Transactions. SSRN 5335578, 2025:\n"
    "Предложен каскад Kernel PCA → LOF → Modern Hopfield / SNN.\n"
    "Недостатки: высокая сложность полной SNN-архитектуры; результаты сложно "
    "воспроизвести без открытого кода на том же датасете.\n\n"
    "Breunig M. et al. LOF: Identifying Density-Based Local Outliers. SIGMOD, 2000:\n"
    "Классический unsupervised-метод по локальной плотности.\n"
    "Недостатки: не использует ассоциативную память нормальных паттернов транзакций.\n\n"
    "Ramsauer H. et al. Hopfield Networks is All You Need. ICLR, 2021:\n"
    "Modern Hopfield Energy для ассоциативного поиска.\n"
    "Недостатки: не адаптирован напрямую к табличным банковским транзакциям.\n\n"
    "computingvictor. Financial Transactions Dataset. Kaggle, 2024:\n"
    "Реальные данные CaixaBank Tech AI Hackathon.\n"
    "Недостатки: экстремальный дисбаланс классов (~0,1% fraud).\n\n"
    "Вывод:\n"
    "Необходима адаптированная и воспроизводимая реализация гибрида LOF + Hopfield "
    "на реальных транзакциях с честной оценкой на сбалансированном тесте."
  ),
  5: (
    "Методы:\n"
    "• Предобработка и feature engineering (StandardScaler, OneHot)\n"
    "• Kernel PCA — нелинейное снижение размерности (RBF-ядро)\n"
    "• LOF — безучительное обнаружение локальных выбросов\n"
    "• Классическая сеть Хопфилда + Modern Hopfield Energy (MHE)\n"
    "• Fusion-скоринг: итоговый скор = 0,35·LOF + 0,65·MHS\n"
    "• Сравнительный эксперимент с supervised и unsupervised baseline\n\n"
    "Инструменты:\n"
    "Python 3.10+, NumPy, Pandas, scikit-learn, Matplotlib, Seaborn, "
    "Jupyter, joblib, kagglehub\n\n"
    "Данные:\n"
    "Kaggle — computingvictor/transactions-fraud-datasets "
    "(CaixaBank, ~13 млн транзакций, ~13 332 метки fraud)"
  ),
  6: (
    "Научная новизна и практическая значимость:\n\n"
    "• Реализован воспроизводимый пайплайн MHSNet (адаптация Zhao et al., SSRN 5335578) "
    "на реальном банковском датасете CaixaBank, а не на синтетических данных.\n\n"
    "• Предложена гибридная схема Kernel PCA → LOF → Hopfield с комбинированным "
    "аномальным скором (ошибка восстановления + расстояние до прототипа + MHE).\n\n"
    "• Память сети Хопфилда обучается только на нормальных транзакциях — метод "
    "применим при дефиците размеченных fraud.\n\n"
    "• Реализована потоковая загрузка 13M+ строк и оценка по метрике Recall@FPR=5%, "
    "релевантной для банковского антифрода.\n\n"
    "• Открытый репозиторий с notebook, кодом и инструкцией воспроизведения эксперимента."
  ),
  7: (
    "Этап 1. Архитектура решения\n"
    "Данные → StandardScaler → Kernel PCA → LOF → Hopfield MHS → fraud / legit\n\n"
    "Этап 2. Датасет CaixaBank (Kaggle)\n"
    "• transactions_data.csv — транзакции\n"
    "• cards_data.csv — атрибуты карт\n"
    "• train_fraud_labels.json — метки fraud\n"
    "• В полной выборке ~13 332 fraud (~0,1%)\n\n"
    "Этап 3. Реализация\n"
    "• 21 признак после объединения и engineering\n"
    "• До 64 прототипов (K-means) в памяти Хопфилда\n"
    "• Подбор порогов LOF и MHS на validation (метрика F1)\n"
    "• 42 763 транзакции в экспериментальной подвыборке"
  ),
  8: (
    "Сбор и подготовка данных\n"
    "• Источник: computingvictor/transactions-fraud-datasets\n"
    "• Потоковое чтение CSV чанками (без загрузки 13M строк в RAM)\n"
    "• Merge транзакций, карт и меток; парсинг суммы, времени, MCC\n"
    "• StandardScaler; сбалансированные val/test (по N на класс)\n\n"
    "Модель MHSNet\n"
    "• Kernel PCA: 15 компонент, RBF-ядро, fit на подвыборке\n"
    "• LOF: k=20, обучение только на нормальных транзакциях\n"
    "• Hopfield: S = 0,35·R̂ + 0,20·D̂ + 0,45·Ê (восстановление, расстояние, MHE)\n"
    "• Fusion: S_final = 0,35·LOF + 0,65·MHS\n\n"
    "Оценка качества\n"
    "• Метрики: accuracy, precision, recall, F1, ROC-AUC, Recall@FPR=5%\n"
    "• Baseline: Logistic Regression, Random Forest, Isolation Forest"
  ),
  9: (
    "Основные результаты (тестовая выборка, CaixaBank):\n\n"
    "Модель                  F1      ROC-AUC   Recall@FPR=5%\n"
    "MHSNet (предложенный)   0,86    0,93      0,69\n"
    "Random Forest           0,94    0,99      0,96\n"
    "Logistic Regression     0,92    0,98      0,91\n"
    "Isolation Forest        0,86    0,92      0,65\n\n"
    "MHSNet: Recall = 0,91; Precision = 0,82; Accuracy = 0,86\n\n"
    "Репозиторий с кодом и экспериментами:\n"
    "hopfield-network-for-fraud-nir\n"
    "• notebooks/MHSNet_Fraud_Detection_NIR.ipynb\n"
    "• README.md — инструкция запуска\n"
    "• src/ — реализация метода"
  ),
  10: (
    "Интерпретация результатов\n"
    "• MHSNet сопоставим с Isolation Forest по F1 (0,86) — оба unsupervised-подхода\n"
    "• Supervised-модели (RF, LR) выше по F1 — ожидаемо при наличии меток при обучении\n"
    "• Recall@FPR=5% = 0,69 — при ≤5% ложных тревог обнаруживается 69% fraud\n"
    "• Modern Hopfield Energy улучшает разделимость vs. только ошибка восстановления\n\n"
    "Выводы\n"
    "• Цель достигнута: метод реализован и проверен на реальных данных\n"
    "• Научная идея MHSNet адаптирована и воспроизведена\n"
    "• Практическая ценность: обучение на нормальных операциях без меток fraud\n\n"
    "Для слайда: вставить графики metrics_comparison.png и roc_curves.png "
    "из папки outputs после запуска run.py"
  ),
  11: (
    "Направления дальнейших исследований:\n\n"
    "• Полная реализация SNN-ветки из оригинальной статьи MHSNet-SNN\n"
    "• Онлайн-обновление прототипов Hopfield (адаптация к concept drift)\n"
    "• Гибридная схема: правила банка + unsupervised MHSNet\n"
    "• Оценка в потоковом режиме, близком к real-time антифроду\n"
    "• Масштабирование Kernel PCA на полной выборке (approximate KPCA)"
  ),
  12: (
    "Заключение\n\n"
    "Выполнены все поставленные задачи:\n"
    "1. Проанализированы методы fraud detection и обнаружения аномалий.\n"
    "2. Подготовлен реальный датасет CaixaBank (Kaggle).\n"
    "3. Реализован пайплайн MHSNet: Kernel PCA → LOF → Hopfield (MHE).\n"
    "4. Проведено сравнение с baseline-моделями.\n"
    "5. Получены метрики F1=0,86, ROC-AUC=0,93, Recall@FPR=5%=0,69.\n\n"
    "Цель работы достигнута: разработан и экспериментально исследован метод "
    "выявления мошеннических транзакций с применением сети Хопфилда.\n\n"
    "Подготовлены: Jupyter notebook для НИР, репозиторий с кодом, "
    "документация docs/METHODOLOGY.md."
  ),
  13: (
    "Степень готовности к публикации:\n\n"
    "• Подготовлен полный текст НИР в формате Jupyter notebook "
    "(теория, эксперименты, выводы, литература).\n"
    "• Оформлена методология с формулами (docs/METHODOLOGY.md).\n"
    "• Репозиторий содержит README с инструкцией воспроизведения эксперимента.\n"
    "• Результаты получены на открытом реальном датасете — работа воспроизводима.\n\n"
    "План публикации:\n"
    "• Оформление тезисов / статьи по результатам НИР "
    "(конференция по ML, анализу данных или финтеху).\n"
    "• Доработка текста с акцентом на адаптацию MHSNet к CaixaBank "
    "и метрику Recall@FPR=5%.\n\n"
    "[Указать: журнал/конференцию, статус подачи — после согласования с руководителем]"
  ),
  14: (
    "Спасибо за внимание!\n\n"
    "Вопросы?\n\n"
    "Литература:\n"
    "1. Zhao Y. MHSNet-SNN. SSRN 5335578, 2025.\n"
    "2. Hopfield J.J. PNAS, 1982.\n"
    "3. Ramsauer H. et al. ICLR, 2021.\n"
    "4. Breunig M. et al. SIGMOD, 2000.\n"
    "5. computingvictor. Kaggle, 2024."
  ),
}


def _set_body(slide, text: str) -> None:
  if len(slide.shapes) == 0:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    box.text_frame.text = text
    return
  if len(slide.shapes) < 3:
    target = slide.shapes[-1]
    if hasattr(target, "text_frame"):
      target.text = text
      return
    raise ValueError(f"Unexpected slide layout with {len(slide.shapes)} shapes")
  body = slide.shapes[2]
  if not hasattr(body, "text_frame"):
    body = slide.shapes[1]
  body.text = text


def _set_title_box(slide, text: str) -> None:
  if slide.shapes and hasattr(slide.shapes[0], "text_frame"):
    slide.shapes[0].text = text


def build_presentation() -> Path:
  if not TEMPLATE.exists():
    raise FileNotFoundError(f"Template not found: {TEMPLATE}")

  OUTPUT.parent.mkdir(parents=True, exist_ok=True)
  shutil.copy2(TEMPLATE, OUTPUT)
  prs = Presentation(str(OUTPUT))

  _set_title_box(prs.slides[0], "Направление подготовки:\nПрикладная математика / Информатика")

  for index, content in SLIDE_CONTENT.items():
    slide = prs.slides[index - 1]
    if index == 1:
      slide.shapes[1].text = content
    else:
      _set_body(slide, content)

  prs.save(str(OUTPUT))
  return OUTPUT


if __name__ == "__main__":
  path = build_presentation()
  print(f"Presentation saved: {path}")
